import sys
from pptx import Presentation

def extract_pptx(pptx_path, out_path):
    prs = Presentation(pptx_path)
    with open(out_path, 'w', encoding='utf-8') as f:
        for i, slide in enumerate(prs.slides):
            f.write(f"=== SLIDE {i+1} ===\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    f.write(shape.text.strip() + "\n")
            f.write("\n")

if __name__ == "__main__":
    extract_pptx("/home/jaimeconde/Descargas/manual_identidad_Ruta_Plata (1).pptx", "manual_identidad.txt")
    print("PPTX content extracted successfully.")
